from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette (matches the dark UI) ──────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # slate-950
PANEL    = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
EMERALD  = RGBColor(0x10, 0xB9, 0x81)   # emerald-500
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)   # amber-400

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl

def box(sl, left, top, w, h, text, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, bg=None, italic=False):
    txBox = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def rect(sl, left, top, w, h, color):
    shape = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_bar(sl):
    rect(sl, 0, 0, 13.33, 0.07, EMERALD)

def bullet_box(sl, left, top, w, h, items, size=18, color=WHITE, spacing=0.55):
    for i, item in enumerate(items):
        box(sl, left, top + i * spacing, w, 0.5, item, size=size, color=color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, BG)
rect(sl, 0, 0, 0.12, 7.5, EMERALD)           # left accent bar
rect(sl, 0, 6.8, 13.33, 0.7, PANEL)          # footer band

box(sl, 0.5, 1.5, 12, 1.2, "VaultAI", size=72, bold=True, color=EMERALD)
box(sl, 0.5, 2.8, 11, 0.7, "The Zero-Knowledge Wealth Copilot", size=28, color=WHITE)
box(sl, 0.5, 3.6, 11, 0.5,
    "A local-first AI agent that manages your taxes, budget & investments",
    size=20, color=MUTED, italic=True)

box(sl, 0.5, 4.5, 5, 0.5, "Team 9  ·  Agentic AI Hackathon", size=16, color=MUTED)
box(sl, 0.5, 6.9, 12, 0.4, "Confidential — Judge Review Copy", size=13,
    color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "The Problem", size=36, bold=True, color=EMERALD)

problems = [
    ("🔒  The Trust Gap",
     "People want AI to manage their money — but fear uploading raw\nbank statements to the cloud. One breach and your financial life is exposed."),
    ("🐟  Goldfish Memory",
     "AI financial bots forget you the moment you close the tab.\nReal planning needs 365 days of context — remembering you had a\nbaby in February to claim a tax credit in April."),
    ("📊  Fragmented Tools",
     "Tax software, budgeting apps, and investment platforms don't\ntalk to each other. You're the integration layer."),
]

for i, (title, body) in enumerate(problems):
    top = 1.3 + i * 1.9
    rect(sl, 0.5, top, 12.3, 1.65, PANEL)
    box(sl, 0.75, top + 0.1, 11.5, 0.45, title, size=20, bold=True, color=EMERALD)
    box(sl, 0.75, top + 0.55, 11.5, 1.0, body, size=16, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "The Solution", size=36, bold=True, color=EMERALD)
box(sl, 0.5, 1.0, 12, 0.45,
    "A Family Office in your pocket — running entirely on your local machine.",
    size=20, color=WHITE, italic=True)

pillars = [
    ("🛡️  Sanitizes Data",
     "Strips PII before the AI sees anything.\nAccount numbers, emails & phone numbers become\nACCT_001, EMAIL_001 — locally, instantly."),
    ("🧠  Builds Memory",
     "Extracts life events from casual chat and persists\nthem to a local vault. 365 days of context,\nnever touching the cloud."),
    ("⚡  Executes Strategy",
     "Runs deterministic cash-flow forecasting, searches\nlive tax codes via Tavily, and recommends\nexact financial actions with sources."),
]

for i, (title, body) in enumerate(pillars):
    left = 0.4 + i * 4.3
    rect(sl, left, 1.7, 4.1, 5.0, PANEL)
    rect(sl, left, 1.7, 4.1, 0.06, EMERALD)
    box(sl, left + 0.2, 1.85, 3.7, 0.55, title, size=19, bold=True, color=EMERALD)
    box(sl, left + 0.2, 2.55, 3.7, 3.8, body, size=16, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "System Architecture", size=36, bold=True, color=EMERALD)

layers = [
    ("Layer 1: PII Masker",       "LOCAL — no network",   EMERALD),
    ("Layer 2: Memory Retrieval", "LOCAL — SQLite + ChromaDB", RGBColor(0x6E, 0xE7, 0xB7)),
    ("Layer 3: LLM Agent",        "CLOUD — masked data only reaches Claude", YELLOW),
    ("Layer 4: Tool Execution",   "LOCAL — forecaster, memory writer (Tavily = exception)", RGBColor(0x6E, 0xE7, 0xB7)),
    ("Layer 5: De-mask & Return", "LOCAL — UIDs swapped back before user sees response", EMERALD),
]

for i, (name, desc, color) in enumerate(layers):
    top = 1.2 + i * 1.1
    rect(sl, 1.5, top, 10.3, 0.85, PANEL)
    rect(sl, 1.5, top, 0.07, 0.85, color)
    box(sl, 1.8, top + 0.05, 5.5, 0.4, name, size=18, bold=True, color=color)
    box(sl, 1.8, top + 0.45, 9.5, 0.35, desc, size=14, color=MUTED)
    if i < len(layers) - 1:
        box(sl, 6.4, top + 0.85, 0.6, 0.25, "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)

box(sl, 0.5, 6.9, 12.3, 0.4,
    "Privacy guarantee: raw financial data is physically incapable of leaving the machine — Docker enforces this at the network layer.",
    size=12, color=MUTED, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "Tech Stack", size=36, bold=True, color=EMERALD)

cols = [
    ("Backend", [
        "Python 3.11 + FastAPI",
        "Anthropic Claude Haiku",
        "Tavily API (web search)",
        "Pandas (forecaster)",
        "Docker + Docker Compose",
    ]),
    ("Memory", [
        "SQLite (life events,",
        "  conversation history)",
        "ChromaDB (vector store)",
        "sentence-transformers",
        "  all-MiniLM-L6-v2",
    ]),
    ("Frontend", [
        "React 19 + Vite",
        "Tailwind CSS v4",
        "React Router v7",
        "3 screens: Chat,",
        "  Upload CSV, Dashboard",
    ]),
    ("Privacy", [
        "Regex PII masker",
        "  (ACCT, EMAIL, PHONE,",
        "  SSN patterns)",
        "UID lookup table",
        "De-mask on response",
    ]),
]

for i, (title, items) in enumerate(cols):
    left = 0.35 + i * 3.2
    rect(sl, left, 1.2, 3.0, 5.8, PANEL)
    rect(sl, left, 1.2, 3.0, 0.06, EMERALD)
    box(sl, left + 0.15, 1.3, 2.7, 0.45, title, size=18, bold=True, color=EMERALD)
    for j, item in enumerate(items):
        box(sl, left + 0.15, 1.95 + j * 0.78, 2.8, 0.6, f"• {item}", size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Demo: Time-Travel
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, 'Demo: "The Time-Travel Scenario"', size=34, bold=True, color=EMERALD)

steps = [
    ("January  📂",  "EMERALD",
     'Upload CSV  →  VaultAI masks PII locally  →  Agent notices heavy tech-gear\nspending and asks: "Are you a freelancer? We can write this off."'),
    ("July  👶",     "YELLOW",
     '"I just had a baby girl!"  →  Agent silently writes\n{life_event: new_child, date: 2026-07} to the local memory vault.'),
    ("Tax Season  💰", "EMERALD",
     '"Prepare my taxes."  →  VaultAI retrieves the July memory, surfaces the\n$2,000 Child Tax Credit, finds $5k surplus, confirms T-Bills at 5% yield.'),
]

for i, (label, col, body) in enumerate(steps):
    top = 1.3 + i * 1.85
    rect(sl, 0.5, top, 12.3, 1.6, PANEL)
    c = EMERALD if col == "EMERALD" else YELLOW
    rect(sl, 0.5, top, 0.07, 1.6, c)
    box(sl, 0.75, top + 0.1, 3.0, 0.45, label, size=20, bold=True, color=c)
    box(sl, 0.75, top + 0.65, 11.5, 0.85, body, size=16, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — What We Built (4-Day Sprint)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "4-Day Build — All Done ✅", size=36, bold=True, color=EMERALD)

days = [
    ("Day 1", "FastAPI backend · Docker · /chat endpoint · /health"),
    ("Day 2", "Tavily web search · Cash-flow forecaster · Claude tool-use loop"),
    ("Day 3", "SQLite + ChromaDB memory · Life events · Conversation history · Semantic recall"),
    ("Day 4", "React frontend · CSV upload · PII masker · De-masking · Docker compose"),
]

for i, (day, desc) in enumerate(days):
    top = 1.3 + i * 1.35
    rect(sl, 0.5, top, 12.3, 1.1, PANEL)
    rect(sl, 0.5, top, 0.07, 1.1, EMERALD)
    box(sl, 0.75, top + 0.1, 1.5, 0.4, day, size=20, bold=True, color=EMERALD)
    box(sl, 2.4,  top + 0.1, 10.0, 0.4, "✅  " + desc, size=18, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Privacy Deep Dive
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "Privacy — Not a Policy, a Guarantee", size=34, bold=True, color=EMERALD)

points = [
    ("🔒  Raw data never leaves your machine",
     "The PII masker has no network access. Docker enforces isolation at the\nnetwork layer — it is physically impossible for raw data to reach the cloud."),
    ("🔑  Stable UID mapping",
     '"john.doe@gmail.com" always becomes EMAIL_001 within a session.\nThe lookup table lives only in RAM — cleared when the server restarts.'),
    ("👁️  What Claude sees",
     "Only masked CSV text + your typed messages. No names, account numbers,\nphone numbers, SSNs, or email addresses ever reach the Anthropic API."),
    ("↩️  De-masking on return",
     "After Claude responds, UIDs are swapped back before you see the answer —\nso you read normal text, not robot codes."),
]

for i, (title, body) in enumerate(points):
    top = 1.2 + i * 1.5
    rect(sl, 0.5, top, 12.3, 1.25, PANEL)
    rect(sl, 0.5, top, 0.07, 1.25, EMERALD)
    box(sl, 0.75, top + 0.08, 11.5, 0.42, title, size=18, bold=True, color=EMERALD)
    box(sl, 0.75, top + 0.55, 11.5, 0.65, body, size=15, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl)
box(sl, 0.5, 0.3, 12, 0.7, "Roadmap to Scale", size=36, bold=True, color=EMERALD)

roadmap = [
    ("Multi-device sync",   "Encrypted cloud backup — server holds blobs it cannot read (client-side encryption)."),
    ("Multi-user",          "Swap SQLite → PostgreSQL (one-line SQLAlchemy config). Add auth layer."),
    ("Better PII masking",  "Replace regex with spaCy NLP model for name/entity detection."),
    ("Zero-knowledge cloud","PII masker runs in sandboxed container — only masked context reaches any server."),
    ("Task execution",      "Agent moves from advice → action: tax filing, investment allocation (regulatory review needed)."),
]

for i, (title, body) in enumerate(roadmap):
    top = 1.25 + i * 1.15
    rect(sl, 0.5, top, 12.3, 0.95, PANEL)
    box(sl, 0.8, top + 0.08, 3.5, 0.38, f"→  {title}", size=17, bold=True, color=EMERALD)
    box(sl, 4.4, top + 0.08, 8.2, 0.75, body, size=15, color=MUTED)

box(sl, 0.5, 7.0, 12.3, 0.35,
    "The local-first architecture is our regulatory moat — we cannot comply with a data request for data we do not hold.",
    size=13, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 0.12, 7.5, EMERALD)
rect(sl, 0, 6.8, 13.33, 0.7, PANEL)

box(sl, 0.5, 1.8, 12, 1.2, "Thank You", size=72, bold=True, color=EMERALD)
box(sl, 0.5, 3.1, 12, 0.6, "Questions?", size=32, color=WHITE)
box(sl, 0.5, 3.9, 12, 0.45,
    "VaultAI — local-first · zero-knowledge · always on your side",
    size=18, color=MUTED, italic=True)

box(sl, 0.5, 5.0, 5.5, 0.4, "Team 9  ·  Agentic AI Hackathon", size=16, color=MUTED)
box(sl, 0.5, 6.9, 12, 0.4, "github.com/Vinoothna99/Team9_AgenticAI_Hackathon",
    size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "VaultAI_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
